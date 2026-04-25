"""Generate AuksionWatch investor pitch deck (PPTX).

Follows the Uzbek startup pitch standard structure (10 slides + extras).
Inspired by `Taqdimot_tayyorlash_O'zb.pptx` guide.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# ───────── DESIGN TOKENS ─────────
NAVY = RGBColor(0x0E, 0x1D, 0x34)
NAVY_DEEP = RGBColor(0x06, 0x10, 0x22)
BLUE = RGBColor(0x00, 0x61, 0xAF)
BLUE_LIGHT = RGBColor(0xE7, 0xF1, 0xFA)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_SOFT = RGBColor(0xFE, 0xF2, 0xF2)
EMERALD = RGBColor(0x04, 0x78, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x4B, 0x55, 0x65)
GREY_LIGHT = RGBColor(0xE2, 0xE7, 0xEF)
PAPER = RGBColor(0xFB, 0xFB, 0xFD)

FONT_BODY = "Inter"  # falls back automatically if not installed
FONT_MONO = "JetBrains Mono"

# 16:9 widescreen
W = Inches(13.333)
H = Inches(7.5)


def add_slide(prs, layout_idx=6):
    """Add blank slide with white-on-navy default."""
    blank = prs.slide_layouts[layout_idx]  # 6 = blank
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    return slide


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill_color, *, line_color=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    return shape


def add_logo_mark(slide, x, y, size=Inches(0.55)):
    """Tiny brand mark — gold square with letter A."""
    sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
    sq.fill.solid()
    sq.fill.fore_color.rgb = BLUE
    sq.line.fill.background()
    tf = sq.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "A"
    run.font.name = FONT_BODY
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_brand_header(slide):
    """Top thin gold strip + logo mark + project name."""
    add_rect(slide, 0, 0, W, Inches(0.06), GOLD)
    add_logo_mark(slide, Inches(0.45), Inches(0.32))
    add_text(slide, Inches(1.1), Inches(0.32), Inches(4), Inches(0.3),
             "AuksionWatch", size=12, bold=True, color=NAVY)
    add_text(slide, Inches(1.1), Inches(0.55), Inches(6), Inches(0.3),
             "E-AUKSION ochiq nazorat tizimi", size=9, color=GREY)


def add_kicker(slide, x, y, text, color=GOLD):
    add_text(slide, x, y, Inches(8), Inches(0.3),
             text.upper(), size=10, bold=True, color=color, font=FONT_MONO)


def add_footer(slide, page_num):
    add_text(slide, Inches(0.45), H - Inches(0.4), Inches(6), Inches(0.3),
             "auksionwatch.uz · CC BY-SA 4.0", size=9, color=GREY)
    add_text(slide, W - Inches(1.5), H - Inches(0.4), Inches(1), Inches(0.3),
             f"{page_num:02d}", size=10, color=GREY,
             align=PP_ALIGN.RIGHT, font=FONT_MONO)


# ───────── SLIDES ─────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ===== SLIDE 1 — Title =====
s = add_slide(prs)
# Navy hero background
add_rect(s, 0, 0, W, H, NAVY)
# Gold accent bar left
add_rect(s, 0, 0, Inches(0.35), H, GOLD)
# Big logo mark center-left
mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.9), Inches(2.6), Inches(1.3), Inches(1.3))
mark.fill.solid(); mark.fill.fore_color.rgb = BLUE
mark.line.fill.background()
tf = mark.text_frame
tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "A"; r.font.size = Pt(72); r.font.bold = True
r.font.color.rgb = WHITE; r.font.name = FONT_BODY

add_text(s, Inches(2.8), Inches(2.55), Inches(10), Inches(0.5),
         "AUKSILKORRUPSIYA HACKATHON · 2026", size=11, bold=True,
         color=GOLD, font=FONT_MONO)
add_text(s, Inches(2.8), Inches(2.95), Inches(10), Inches(1.2),
         "AuksionWatch", size=58, bold=True, color=WHITE)
add_text(s, Inches(2.8), Inches(4.0), Inches(10), Inches(0.6),
         "E-AUKSION ochiq nazorat tizimi", size=20, color=GREY_LIGHT)
add_text(s, Inches(2.8), Inches(4.55), Inches(10), Inches(0.5),
         "Davlat mol-mulki auksionlaridagi shubhali sxemalarni AI yordamida aniqlash",
         size=14, color=RGBColor(0xA1, 0xB1, 0xC9))

# Bottom contact strip
add_rect(s, 0, H - Inches(0.9), W, Inches(0.9), NAVY_DEEP)
add_text(s, Inches(0.7), H - Inches(0.7), Inches(6), Inches(0.3),
         "auksionwatch.uz · github.com/.../aksilkorupsiya2026", size=11,
         color=GOLD, font=FONT_MONO)
add_text(s, W - Inches(5), H - Inches(0.7), Inches(4.5), Inches(0.3),
         "Toshkent · 2026", size=11, color=GREY_LIGHT, align=PP_ALIGN.RIGHT,
         font=FONT_MONO)


# ===== SLIDE 2 — Problem =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "MUAMMO · 1/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "130 mlrd so'mlik yer 120 mlrd ga sotildi.", size=36, bold=True, color=NAVY)
add_text(s, Inches(0.45), Inches(2.7), Inches(11), Inches(0.6),
         "Va bu — mingdan biri.", size=22, color=RED, bold=True)

# 3 stat cards
stats = [
    ("130", "mlrd so'm", "Ortiqov keysida\ndavlat zarari", RED),
    ("4.2", "trln so'm", "2026-yanvar prezident\nbayonotidagi yillik zarar", RED),
    ("9 266", "ta lot", "Faqat Fargʻonada — 81% i\n3 odam qo'lida", GOLD),
]
x = Inches(0.45)
for i, (num, unit, label, c) in enumerate(stats):
    box_x = x + Inches(i * 4.15)
    add_rect(s, box_x, Inches(3.7), Inches(3.95), Inches(2.4), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_rect(s, box_x, Inches(3.7), Inches(0.1), Inches(2.4), c)
    add_text(s, box_x + Inches(0.3), Inches(3.85), Inches(3.5), Inches(0.7),
             num, size=42, bold=True, color=c)
    add_text(s, box_x + Inches(0.3), Inches(4.55), Inches(3.5), Inches(0.4),
             unit, size=14, color=NAVY, bold=True)
    add_text(s, box_x + Inches(0.3), Inches(5.0), Inches(3.5), Inches(1),
             label, size=11, color=GREY)

add_text(s, Inches(0.45), Inches(6.5), Inches(12), Inches(0.5),
         "OECD: davlat loyihalari qiymatining 20–30% i korrupsiyaga yo'qoladi.",
         size=12, color=GREY, font=FONT_MONO)
add_footer(s, 2)


# ===== SLIDE 3 — Solution =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "YECHIM · 2/10", color=BLUE)
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "Civic AI — har bir auksion ochiq nazoratda.", size=32, bold=True, color=NAVY)
add_text(s, Inches(0.45), Inches(2.7), Inches(12), Inches(1),
         "AuksionWatch e-auksion.uz dagi 23M+ lotni xalqaro standartlar (OECD, UNCAC, "
         "OCDS, FATF) asosida AI yordamida tahlil qilib, korrupsion sxemalarni "
         "ОLDINDAN aniqlaydi va fuqaroga, jurnalistga, NGO ga ochiq taqdim etadi.",
         size=14, color=GREY)

# 3-pillar diagram
pillars = [
    ("Rule Engine", "OECD/Fazekas\n18+ rule, izohlanadi", BLUE),
    ("ML Ensemble", "XGBoost + IsoForest\nCV AUC 0.98", GOLD),
    ("PEP Layer", "FATF R12\nMansabdor screening", RED),
]
for i, (title, desc, c) in enumerate(pillars):
    bx = Inches(0.45 + i * 4.15)
    add_rect(s, bx, Inches(4.0), Inches(3.95), Inches(2.0), WHITE,
             line_color=c, radius=True)
    add_rect(s, bx + Inches(0.3), Inches(4.2), Inches(0.5), Inches(0.5), c, radius=True)
    add_text(s, bx + Inches(0.95), Inches(4.25), Inches(3), Inches(0.5),
             title, size=15, bold=True, color=NAVY)
    add_text(s, bx + Inches(0.3), Inches(5.0), Inches(3.5), Inches(1),
             desc, size=11, color=GREY)

add_rect(s, Inches(0.45), Inches(6.3), Inches(12.4), Inches(0.7), BLUE_LIGHT, radius=True)
add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
         "→ Web dashboard · Telegram bot · Public REST API + CSV/JSON eksport",
         size=12, color=BLUE, bold=True, font=FONT_MONO)
add_footer(s, 3)


# ===== SLIDE 4 — Market =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "BOZOR · 3/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "Bozor o'lchami va auditoriya", size=32, bold=True, color=NAVY)

add_text(s, Inches(0.45), Inches(2.7), Inches(8), Inches(0.5),
         "TAM: e-auksion yillik aylanma · 2023:", size=12, color=GREY)
add_text(s, Inches(0.45), Inches(3.1), Inches(8), Inches(0.8),
         "19.7 trln so'm", size=42, bold=True, color=BLUE)

# Market segments
segs = [
    ("Davlat organlari", "B2G", "Aksilkorrupsiya agentligi,\nHisob Palatasi, Davaktiv",
     "$50–500K/yil"),
    ("Mediya / NGO", "B2B", "Kun.uz, Gazeta.uz, Spot.uz,\nCABAR Asia, OCCRP",
     "$200–1K/oy"),
    ("Donor grant", "Grant", "UNDP, USAID, EU,\nWorld Bank, Soros",
     "$50K–2M"),
    ("Korxona compliance", "B2B", "Bank due-diligence,\nyuridik firmalar, FDI",
     "$500/check"),
]
y0 = Inches(4.3)
for i, (name, kind, desc, price) in enumerate(segs):
    bx = Inches(0.45 + (i % 2) * 6.4)
    by = y0 + Inches((i // 2) * 1.4)
    add_rect(s, bx, by, Inches(6.2), Inches(1.25), WHITE, line_color=GREY_LIGHT, radius=True)
    add_rect(s, bx, by, Inches(0.08), Inches(1.25), GOLD)
    add_text(s, bx + Inches(0.25), by + Inches(0.1), Inches(4), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s, bx + Inches(4.5), by + Inches(0.1), Inches(1.5), Inches(0.4),
             kind, size=10, bold=True, color=BLUE, font=FONT_MONO,
             align=PP_ALIGN.RIGHT)
    add_text(s, bx + Inches(0.25), by + Inches(0.5), Inches(5.5), Inches(0.7),
             desc, size=10, color=GREY)
    add_text(s, bx + Inches(0.25), by + Inches(0.95), Inches(5.5), Inches(0.3),
             price, size=10, color=GOLD, bold=True, font=FONT_MONO)

add_footer(s, 4)


# ===== SLIDE 5 — Competitors =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "RAQOBATCHILAR · 4/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(12), Inches(1),
         "Bizning noyob pozitsiya", size=32, bold=True, color=NAVY)

# Comparison table
rows = [
    ("Funksiya",          "e-anticor.uz",    "data.egov.uz",   "AuksionWatch"),
    ("Lot tahlili",       "Ichki",           "Yo'q",           "✓ 11K+ live"),
    ("Public API",        "Yo'q",            "Cheklangan",     "✓ Swagger + CSV"),
    ("AI / ML",           "Yopiq tizim",     "Yo'q",           "✓ XGB AUC 0.98"),
    ("Telegram bot",      "Yo'q",            "Yo'q",           "✓ /check, /firma"),
    ("PEP screening",     "Ichki",           "Yo'q",           "✓ FATF R12"),
    ("Open source",       "✗",               "✗",              "✓ CC BY-SA 4.0"),
    ("Foydalanuvchi",     "Davlat",          "Davlat",         "Fuqaro · Jurnalist"),
]
col_w = [Inches(3.0), Inches(2.7), Inches(2.7), Inches(3.0)]
x_start = Inches(0.85)
y_start = Inches(2.9)
row_h = Inches(0.42)
for ri, row in enumerate(rows):
    is_header = ri == 0
    for ci, cell in enumerate(row):
        cx = x_start + sum(col_w[:ci], Emu(0))
        cy = y_start + row_h * ri
        is_us = ci == 3
        bg = NAVY if is_header else (BLUE_LIGHT if is_us and not is_header else WHITE)
        if is_header:
            add_rect(s, cx, cy, col_w[ci], row_h, NAVY)
        elif is_us:
            add_rect(s, cx, cy, col_w[ci], row_h, BLUE_LIGHT)
        else:
            add_rect(s, cx, cy, col_w[ci], row_h, WHITE, line_color=GREY_LIGHT)
        text_color = WHITE if is_header else (BLUE if is_us else NAVY)
        bold = is_header or is_us or ci == 0
        add_text(s, cx + Inches(0.15), cy + Inches(0.06), col_w[ci] - Inches(0.2),
                 row_h, cell, size=11, bold=bold, color=text_color,
                 font=FONT_MONO if is_header else FONT_BODY)

add_footer(s, 5)


# ===== SLIDE 6 — Business model =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "BIZNES MODEL · 5/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "Daromadning 4 oqimi", size=32, bold=True, color=NAVY)

# Revenue streams
streams = [
    ("01", "Davlat kontraktlari",
     "Aksilkorrupsiya agentligi va Hisob Palatasi uchun audit asbob.\n"
     "Ortikhov keysidan keyin ehtiyoj kuchli.", "B2G"),
    ("02", "Donor grantlari",
     "UNDP, USAID, EU, World Bank — anti-corruption / open data programs.\n"
     "Ukraine ProZorro $5M+ olgan.", "Grant"),
    ("03", "Media / NGO subscription",
     "Jurnalist va tergov NGO uchun premium API + alert.\n"
     "Hozir bepul, kelgusida $200–1K/oy.", "B2B"),
    ("04", "Compliance check",
     "Bank, yuridik firma, FDI investor uchun firma tekshirish.\n"
     "API kalit, $500/raport.", "B2B"),
]
y0 = Inches(2.9)
for i, (num, title, desc, kind) in enumerate(streams):
    by = y0 + Inches(i * 0.95)
    add_rect(s, Inches(0.45), by, Inches(12.4), Inches(0.85), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_text(s, Inches(0.65), by + Inches(0.18), Inches(0.7), Inches(0.5),
             num, size=24, bold=True, color=GOLD, font=FONT_MONO)
    add_text(s, Inches(1.5), by + Inches(0.05), Inches(8), Inches(0.4),
             title, size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.5), by + Inches(0.4), Inches(8.5), Inches(0.5),
             desc, size=10, color=GREY)
    add_rect(s, Inches(11.5), by + Inches(0.25), Inches(1.3), Inches(0.35),
             BLUE_LIGHT, radius=True)
    add_text(s, Inches(11.5), by + Inches(0.27), Inches(1.3), Inches(0.3),
             kind, size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER,
             font=FONT_MONO)
add_footer(s, 6)


# ===== SLIDE 7 — Status / Traction =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "HOLAT · 6/10", color=EMERALD)
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "MVP live · 11,204 lot tahlilda", size=32, bold=True, color=NAVY)

# Traction grid
grid = [
    ("11,204", "Real lot DB'da", BLUE),
    ("6,784", "Yuqori xavfli aniqlandi", RED),
    ("186", "ML KRITIK (top 2%)", GOLD),
    ("9 266", "Fargʻona Excel ingest", BLUE),
    ("1 938", "API real-time scrape", BLUE),
    ("18+", "Active red-flag rule", GOLD),
    ("0.98", "XGBoost CV AUC", EMERALD),
    ("8", "PEP watchlist'da", RED),
]
for i, (num, label, c) in enumerate(grid):
    bx = Inches(0.45 + (i % 4) * 3.15)
    by = Inches(2.9 + (i // 4) * 1.4)
    add_rect(s, bx, by, Inches(2.95), Inches(1.2), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_rect(s, bx, by, Inches(0.06), Inches(1.2), c)
    add_text(s, bx + Inches(0.2), by + Inches(0.1), Inches(2.7), Inches(0.65),
             num, size=24, bold=True, color=c, font=FONT_MONO)
    add_text(s, bx + Inches(0.2), by + Inches(0.75), Inches(2.7), Inches(0.4),
             label, size=10, color=GREY)

# Live URLs
add_rect(s, Inches(0.45), Inches(5.95), Inches(12.4), Inches(1), NAVY, radius=True)
add_text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.4),
         "LIVE PRODUCTION", size=10, bold=True, color=GOLD, font=FONT_MONO)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "aksilkorrupsiya2026.up.railway.app  ·  Railway: backend + Postgres + frontend deployed",
         size=12, color=WHITE, font=FONT_MONO)
add_footer(s, 7)


# ===== SLIDE 8 — Methodology / Tech stack =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "METODOLOGIYA · 7/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "5 OECD/OCP toifasi · 2 qatlamli AI", size=32, bold=True, color=NAVY)

# 5 categories
cats = [
    ("A", "Past shaffoflik", "yopiq auksion, qisqa muddat", RGBColor(0xC2, 0x41, 0x0C)),
    ("B", "Kelishuv", "monopol sotuvchi, PEP", RGBColor(0x7E, 0x22, 0xCE)),
    ("C", "Bid-rigging", "1 ishtirokchi, takror", RED),
    ("D", "Firibgarlik", "past baho, diskont", RGBColor(0xCA, 0x8A, 0x04)),
    ("E", "Past raqobat", "regional dominance", RGBColor(0x08, 0x91, 0xB2)),
]
for i, (code, name, desc, c) in enumerate(cats):
    bx = Inches(0.45 + i * 2.55)
    add_rect(s, bx, Inches(2.9), Inches(2.4), Inches(2.0), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_rect(s, bx + Inches(0.2), Inches(3.1), Inches(0.5), Inches(0.5), c, radius=True)
    add_text(s, bx + Inches(0.2), Inches(3.18), Inches(0.5), Inches(0.4),
             code, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(0.2), Inches(3.7), Inches(2.1), Inches(0.5),
             name, size=12, bold=True, color=NAVY)
    add_text(s, bx + Inches(0.2), Inches(4.15), Inches(2.1), Inches(0.7),
             desc, size=9, color=GREY)

# Stack row
stack_items = [
    "Python · FastAPI",
    "Next.js · React 19",
    "PostgreSQL · SQLite",
    "XGBoost · Scikit-learn",
    "Playwright · httpx",
    "Tailwind · Recharts",
    "Telegram · aiogram 3",
    "Railway · GitHub",
]
y_stack = Inches(5.4)
add_text(s, Inches(0.45), y_stack, Inches(8), Inches(0.3),
         "TECHNOLOGY", size=10, bold=True, color=GOLD, font=FONT_MONO)
for i, t in enumerate(stack_items):
    bx = Inches(0.45 + (i % 4) * 3.15)
    by = y_stack + Inches(0.4 + (i // 4) * 0.55)
    add_rect(s, bx, by, Inches(2.95), Inches(0.45), BLUE_LIGHT, radius=True)
    add_text(s, bx + Inches(0.15), by + Inches(0.08), Inches(2.7), Inches(0.3),
             t, size=10, color=BLUE, font=FONT_MONO, bold=True)
add_footer(s, 8)


# ===== SLIDE 9 — Team =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "JAMOA · 8/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "4 mutaxassis · 8 soat MVP", size=32, bold=True, color=NAVY)

team = [
    ("Backend Engineer", "FastAPI · Postgres · Risk engine",
     "11,204 lot ingest, 18+ rule, OECD taksonomiya", BLUE),
    ("Frontend Engineer", "Next.js · Tailwind · Leaflet",
     "8 sahifa, dark mode, PDF export, my.gov.uz dizayn", GOLD),
    ("ML Engineer", "XGBoost · IsolationForest · Pandas",
     "33 feature, CV AUC 0.98, weak-supervised pipeline", RED),
    ("Designer / PM", "Pitch · UX · Methodology",
     "5 toifa OECD, chetel research, demo skripti", EMERALD),
]
for i, (role, stack, contrib, c) in enumerate(team):
    bx = Inches(0.45 + (i % 2) * 6.4)
    by = Inches(2.9 + (i // 2) * 1.85)
    add_rect(s, bx, by, Inches(6.2), Inches(1.7), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_rect(s, bx, by, Inches(0.1), Inches(1.7), c)
    # avatar circle
    add_rect(s, bx + Inches(0.3), by + Inches(0.3), Inches(0.8), Inches(0.8),
             c, radius=True)
    add_text(s, bx + Inches(0.3), by + Inches(0.45), Inches(0.8), Inches(0.5),
             role[:1], size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx + Inches(1.3), by + Inches(0.25), Inches(4.7), Inches(0.4),
             role, size=14, bold=True, color=NAVY)
    add_text(s, bx + Inches(1.3), by + Inches(0.65), Inches(4.7), Inches(0.4),
             stack, size=10, color=BLUE, font=FONT_MONO)
    add_text(s, bx + Inches(1.3), by + Inches(1.05), Inches(4.7), Inches(0.6),
             contrib, size=10, color=GREY)
add_footer(s, 9)


# ===== SLIDE 10 — Roadmap & ask =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "REJA · 9/10")
add_text(s, Inches(0.45), Inches(1.7), Inches(11), Inches(1),
         "Yo'l xaritasi va so'rov", size=32, bold=True, color=NAVY)

# Roadmap
phases = [
    ("Q2 2026",   "MVP",          "11K lot, 5 toifa, ML ensemble, Railway production", GOLD),
    ("Q3 2026",   "v2 Scale",     "5K → 100K lot · NetworkX graf · OpenSanctions PEP API", BLUE),
    ("Q4 2026",   "Cross-platform", "xarid.uzex · regulation.gov · decisions.esud integratsiya", BLUE),
    ("Q1 2027",   "v3 Mobile",    "iOS/Android ilova · push alert · whistleblower form", BLUE),
    ("2027+",     "Halqaro",      "OCDS Cardinal hamkorligi · UN SDG 16 · EU GPA", EMERALD),
]
for i, (q, name, desc, c) in enumerate(phases):
    by = Inches(2.85 + i * 0.65)
    add_rect(s, Inches(0.45), by, Inches(7.5), Inches(0.55), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_text(s, Inches(0.65), by + Inches(0.1), Inches(1.3), Inches(0.4),
             q, size=11, bold=True, color=c, font=FONT_MONO)
    add_text(s, Inches(2.0), by + Inches(0.1), Inches(1.5), Inches(0.4),
             name, size=12, bold=True, color=NAVY)
    add_text(s, Inches(3.5), by + Inches(0.13), Inches(4.4), Inches(0.4),
             desc, size=10, color=GREY)

# Ask block
add_rect(s, Inches(8.3), Inches(2.85), Inches(4.6), Inches(3.8), NAVY, radius=True)
add_text(s, Inches(8.55), Inches(3.0), Inches(4), Inches(0.4),
         "BIZGA KERAK", size=10, bold=True, color=GOLD, font=FONT_MONO)
asks = [
    "$50K seed grant",
    "OpenSanctions API kreditlari",
    "Davlat hamkorlik (NDA)",
    "OCCRP / CABAR mentor",
    "Server resurs (Railway / VPS)",
]
for i, a in enumerate(asks):
    add_text(s, Inches(8.7), Inches(3.5 + i * 0.55), Inches(4), Inches(0.4),
             "→", size=14, color=GOLD, bold=True, font=FONT_MONO)
    add_text(s, Inches(9.05), Inches(3.5 + i * 0.55), Inches(4), Inches(0.4),
             a, size=12, color=WHITE)

add_footer(s, 10)


# ===== SLIDE 11 — International precedents =====
s = add_slide(prs)
add_brand_header(s)
add_kicker(s, Inches(0.45), Inches(1.3), "DELILLAR · 10/10", color=EMERALD)
add_text(s, Inches(0.45), Inches(1.7), Inches(12), Inches(1),
         "Xalqaro tajriba: bu yo'l ishlaydi.", size=32, bold=True, color=NAVY)

cases = [
    ("🇺🇦", "Ukraine ProZorro + DOZORRO",
     "Open-source civic AI, 2014-dan beri. Bizning eng yaqin model.", "$6 mlrd"),
    ("🇧🇷", "Brazil ALICE (TCU)",
     "AI tender risk pre-detection, 2017-dan. 139,566 tender/yil tahlil.", "EUR 35M+/yil"),
    ("🇨🇿", "Czech zIndex.cz",
     "Hokimliklar shaffoflik reytingi. Yuqori reytingli organlar har tenderda 5% kam to'laydi.", "−5% xarajat"),
    ("🇸🇰", "Slovakia Open Tender (2011)",
     "Barcha shartnomalar majburiy ravishda ochiq. 2% → 50% elektron, +1 ishtirokchi.", "−2-3% narx"),
]
for i, (flag, name, desc, num) in enumerate(cases):
    by = Inches(2.9 + i * 0.95)
    add_rect(s, Inches(0.45), by, Inches(12.4), Inches(0.85), WHITE,
             line_color=GREY_LIGHT, radius=True)
    add_text(s, Inches(0.7), by + Inches(0.18), Inches(0.7), Inches(0.5),
             flag, size=22)
    add_text(s, Inches(1.4), by + Inches(0.05), Inches(7.5), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.4), by + Inches(0.4), Inches(8.5), Inches(0.5),
             desc, size=10, color=GREY)
    add_text(s, Inches(10.3), by + Inches(0.22), Inches(2.4), Inches(0.5),
             num, size=18, bold=True, color=EMERALD, font=FONT_MONO,
             align=PP_ALIGN.RIGHT)

add_text(s, Inches(0.45), Inches(6.85), Inches(12), Inches(0.4),
         "Akademik asos: Mihály Fazekas (Cambridge/ERCAS) — Single Bidder Indicator (CRI),",
         size=10, color=GREY)
add_text(s, Inches(0.45), Inches(7.05), Inches(12), Inches(0.4),
         "OECD Anti-Corruption Outlook 2026, World Bank IACRC, Open Contracting Cardinal",
         size=10, color=GREY)
add_footer(s, 11)


# ===== SLIDE 12 — Contact =====
s = add_slide(prs)
add_rect(s, 0, 0, W, H, NAVY)
add_rect(s, 0, 0, Inches(0.35), H, GOLD)
# big mark
mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
mark.fill.solid(); mark.fill.fore_color.rgb = BLUE; mark.line.fill.background()
tf = mark.text_frame
tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "A"; r.font.size = Pt(80); r.font.bold = True
r.font.color.rgb = WHITE

add_text(s, 0, Inches(3.0), W, Inches(0.6),
         "Rahmat!", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.85), W, Inches(0.5),
         "Har bir auksion ochiq nazoratda bo'lsin.",
         size=18, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, 0, Inches(4.6), W, Inches(0.5),
         "aksilkorrupsiya2026.up.railway.app", size=16, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
add_text(s, 0, Inches(5.0), W, Inches(0.5),
         "github.com/testuchunmaxsus/aksilkorupsiya2026", size=12,
         color=GREY_LIGHT, align=PP_ALIGN.CENTER, font=FONT_MONO)
add_text(s, 0, Inches(5.4), W, Inches(0.5),
         "Telegram: @AuksionWatch_bot", size=12,
         color=GREY_LIGHT, align=PP_ALIGN.CENTER, font=FONT_MONO)

# Footer band
add_rect(s, 0, H - Inches(1), W, Inches(1), NAVY_DEEP)
add_text(s, 0, H - Inches(0.7), W, Inches(0.4),
         "OECD · UNCAC · OCDS · FATF R12  ·  CC BY-SA 4.0  ·  Toshkent 2026",
         size=11, color=GOLD, align=PP_ALIGN.CENTER, font=FONT_MONO)


# ───────── SAVE ─────────
out = Path("docs/AuksionWatch_pitch.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Saved {out} ({out.stat().st_size//1024} KB · {len(prs.slides)} slides)")
